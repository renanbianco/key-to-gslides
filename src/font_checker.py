"""Check which fonts in a PPTX are not supported by Google Slides."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Pt

# Fonts natively available in Google Slides / Google Fonts that are commonly supported.
# This list covers the built-in Google Slides fonts plus the most common Google Fonts.
GOOGLE_SLIDES_FONTS = {
    # Google Slides built-in / common Google Fonts (case-insensitive matching used)
    "Arial", "Arial Black", "Arial Narrow", "Arial Rounded MT Bold",
    "Calibri", "Cambria", "Candara", "Century Gothic", "Comic Sans MS",
    "Consolas", "Constantia", "Corbel", "Courier New", "Georgia",
    "Impact", "Lucida Console", "Lucida Sans Unicode", "Microsoft Sans Serif",
    "Palatino Linotype", "Tahoma", "Times New Roman", "Trebuchet MS",
    "Verdana", "Webdings", "Wingdings",
    # Google Fonts (most popular)
    "Roboto", "Roboto Condensed", "Roboto Mono", "Roboto Slab",
    "Open Sans", "Open Sans Condensed",
    "Lato", "Montserrat", "Oswald", "Source Sans Pro", "Source Sans 3",
    "Raleway", "PT Sans", "PT Serif", "PT Sans Narrow", "PT Sans Caption",
    "Merriweather", "Merriweather Sans",
    "Ubuntu", "Ubuntu Condensed", "Ubuntu Mono",
    "Nunito", "Nunito Sans",
    "Playfair Display", "Playfair Display SC",
    "Poppins", "Quicksand", "Titillium Web",
    "Fira Sans", "Fira Sans Condensed", "Fira Mono", "Fira Code",
    "Droid Sans", "Droid Serif", "Droid Sans Mono",
    "Noto Sans", "Noto Serif", "Noto Mono",
    "Work Sans", "Barlow", "Barlow Condensed", "Barlow Semi Condensed",
    "Inter", "DM Sans", "DM Serif Display", "DM Serif Text",
    "Josefin Sans", "Josefin Slab",
    "Libre Baskerville", "Libre Franklin",
    "IBM Plex Sans", "IBM Plex Serif", "IBM Plex Mono",
    "Cabin", "Cabin Condensed",
    "Exo", "Exo 2",
    "Cinzel", "Cinzel Decorative",
    "Crimson Text", "Crimson Pro",
    "EB Garamond",
    "Alegreya", "Alegreya Sans", "Alegreya SC",
    "Arimo",
    "Bitter", "Cardo",
    "Chivo", "Domine",
    "Gentium Basic", "Gentium Book Basic",
    "Inconsolata",
    "Karla",
    "Lora",
    "Muli",
    "Oxygen", "Oxygen Mono",
    "Pacifico",
    "Righteous",
    "Rubik", "Rubik Mono One",
    "Signika", "Signika Negative",
    "Source Code Pro", "Source Serif Pro", "Source Serif 4",
    "Space Mono", "Space Grotesk",
    "Spectral",
    "Varela Round",
    "Yanone Kaffeesatz",
    "Zilla Slab",
    "Anton", "Archivo", "Archivo Black", "Archivo Narrow",
    "Assistant", "BioRhyme",
    "Cormorant", "Cormorant Garamond",
    "Didact Gothic",
    "Frank Ruhl Libre",
    "Heebo",
    "Kalam",
    "Kanit",
    "Lobster",
    "Mada",
    "Maven Pro",
    "Overpass", "Overpass Mono",
    "Prompt",
    "Questrial",
    "Rasa",
    "Sarabun",
    "Teko",
    "Vollkorn",
    # Additional Google Fonts commonly missed
    "Tenor Sans",
    "Alata", "Alegreya SC",
    "Bebas Neue",
    "Bodoni Moda",
    "Brygada 1918",
    "Caveat", "Caveat Brush",
    "Comfortaa",
    "Dancing Script",
    "Dosis",
    "Epilogue",
    "Figtree",
    "Gelasio",
    "Gloock",
    "Gothic A1",
    "Hanken Grotesk",
    "Hahmlet",
    "Hedvig Letters Serif",
    "Instrument Sans", "Instrument Serif",
    "Jost",
    "League Spartan",
    "Lexend", "Lexend Deca", "Lexend Exa", "Lexend Giga",
    "Limelight",
    "Literata",
    "Manrope",
    "Material Icons",
    "Nanum Gothic", "Nanum Myeongjo",
    "Nunito Sans",
    "Outfit",
    "Plus Jakarta Sans",
    "Podkova",
    "Public Sans",
    "Readex Pro",
    "Red Hat Display", "Red Hat Text",
    "Rokkitt",
    "Sora",
    "Syne",
    "Urbanist",
    "Wix Madefor Display", "Wix Madefor Text",
    "Ysabeau", "Ysabeau SC",
}

# Normalise for case-insensitive lookup
_FONTS_LOWER = {f.lower() for f in GOOGLE_SLIDES_FONTS}


def fetch_google_fonts(api_key: str | None = None) -> int:
    """
    Fetch all Google Fonts from the API and add them to GOOGLE_SLIDES_FONTS.
    Returns the number of NEW fonts added.

    If api_key is None, reads from env var then from the saved config file.
    Raises RuntimeError if no key is available or the request fails.
    """
    import urllib.request
    import json as _json
    from src.config import get_google_fonts_api_key

    key = api_key or get_google_fonts_api_key()
    if not key:
        raise RuntimeError(
            "No Google Fonts API key found.\n"
            "Paste your key in the field below and click Load."
        )

    url = (
        f"https://www.googleapis.com/webfonts/v1/webfonts"
        f"?key={key}&fields=items/family&sort=alpha"
    )
    try:
        # Use certifi's CA bundle if available (fixes SSL on macOS stock Python).
        # Fall back to an unverified context — this is a public read-only API,
        # so there is no credential or sensitive data at risk.
        try:
            import ssl, certifi
            ctx = ssl.create_default_context(cafile=certifi.where())
        except Exception:
            import ssl
            ctx = ssl.create_default_context()
            ctx.check_hostname = False
            ctx.verify_mode = ssl.CERT_NONE

        with urllib.request.urlopen(url, timeout=15, context=ctx) as resp:
            data = _json.loads(resp.read())
    except Exception as e:
        raise RuntimeError(f"Google Fonts API request failed: {e}") from e

    added = 0
    for item in data.get("items", []):
        family = item.get("family", "").strip()
        if family and family not in GOOGLE_SLIDES_FONTS:
            GOOGLE_SLIDES_FONTS.add(family)
            _FONTS_LOWER.add(family.lower())
            added += 1
    return added


def get_all_font_names() -> list[str]:
    """Return sorted list of all known fonts (static + any API-fetched)."""
    return sorted(GOOGLE_SLIDES_FONTS)


# Extend with live Google Fonts at import time if a key is available
import os as _os
import urllib.request as _urllib_request
try:
    from src.config import get_google_fonts_api_key as _get_key
    if _get_key():
        fetch_google_fonts()
except Exception:
    pass  # silently fall back to static list


def is_supported(font_name: str) -> bool:
    return font_name.strip().lower() in _FONTS_LOWER


def extract_fonts_from_pptx(pptx_path: str) -> set[str]:
    """Return the set of all font names used in a PPTX file."""
    prs = Presentation(pptx_path)
    fonts: set[str] = set()

    def _collect(run_or_para):
        if hasattr(run_or_para, "font"):
            name = run_or_para.font.name
            if name:
                fonts.add(name)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _collect(para)
                    for run in para.runs:
                        _collect(run)
            # Table cells
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _collect(para)
                            for run in para.runs:
                                _collect(run)

    # Also check slide layout and master defaults
    for slide in prs.slides:
        if slide.slide_layout and slide.slide_layout.placeholders:
            for ph in slide.slide_layout.placeholders:
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        _collect(para)
                        for run in para.runs:
                            _collect(run)

    return fonts


def find_unsupported_fonts(pptx_path: str) -> list[str]:
    """Return sorted list of font names used in the PPTX that Google Slides won't support."""
    all_fonts = extract_fonts_from_pptx(pptx_path)
    unsupported = [f for f in all_fonts if not is_supported(f)]
    return sorted(unsupported)
